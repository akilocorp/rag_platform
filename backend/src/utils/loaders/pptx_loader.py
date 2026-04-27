"""
Minimal PowerPoint (.pptx) loader using python-pptx.

Walks each slide, joins text from every shape that exposes a text frame
(titles, bullets, text boxes), and returns one LangChain Document per slide
with metadata {slide_number, source}. Tables and notes are not included in
this v1 — text-only.
"""
from langchain_core.documents import Document

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - import error surfaced at first use
    Presentation = None


class SimplePPTXLoader:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        if Presentation is None:
            raise ImportError("python-pptx is required to load .pptx files")

        prs = Presentation(self.file_path)
        documents = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in paragraph.runs).strip()
                    if text:
                        lines.append(text)

            content = '\n'.join(lines).strip()
            if not content:
                continue

            documents.append(Document(
                page_content=content,
                metadata={
                    'slide_number': idx,
                    'source': self.file_path,
                },
            ))
        return documents
